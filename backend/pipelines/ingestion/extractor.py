"""
Extract text + structural metadata from various file formats.
Returns a list of PageSection objects preserving heading hierarchy and page numbers.
"""

from dataclasses import dataclass, field
from pathlib import Path


@dataclass
class PageSection:
    text: str
    section: str = ""       # heading path e.g. "Specifications / Power"
    page_number: int | None = None
    has_table: bool = False
    sheet_name: str = ""    # for Excel


def extract(file_path: str) -> list[PageSection]:
    path = Path(file_path)
    ext = path.suffix.lstrip(".").lower()

    extractors = {
        "pdf": _extract_pdf,
        "docx": _extract_docx,
        "xlsx": _extract_xlsx,
        "pptx": _extract_pptx,
        "md": _extract_text,
        "txt": _extract_text,
    }

    fn = extractors.get(ext)
    if fn is None:
        raise ValueError(f"Unsupported file type: {ext}")
    return fn(file_path)


def _extract_pdf(file_path: str) -> list[PageSection]:
    sections = []
    try:
        import pdfplumber
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                text = page.extract_text() or ""
                tables = page.extract_tables()
                has_table = bool(tables)

                # Convert tables to readable text and append
                table_text = ""
                for table in tables:
                    for row in table:
                        if row:
                            table_text += " | ".join(str(c or "") for c in row) + "\n"

                full_text = text
                if table_text:
                    full_text += "\n[TABLE]\n" + table_text

                if full_text.strip():
                    sections.append(PageSection(
                        text=full_text.strip(),
                        page_number=page_num,
                        has_table=has_table,
                    ))
    except Exception:
        # Fallback to pypdf
        import pypdf
        reader = pypdf.PdfReader(file_path)
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                sections.append(PageSection(text=text.strip(), page_number=i + 1))

    # Enrich with heading paths
    _enrich_heading_paths(sections)
    return sections


def _extract_docx(file_path: str) -> list[PageSection]:
    from docx import Document
    doc = Document(file_path)

    sections = []
    heading_stack: list[str] = []
    current_lines: list[str] = []
    current_section = ""
    has_table = False

    def flush():
        nonlocal current_lines, has_table
        text = "\n".join(current_lines).strip()
        if text:
            sections.append(PageSection(
                text=text,
                section=current_section,
                has_table=has_table,
            ))
        current_lines = []
        has_table = False

    for para in doc.paragraphs:
        style = para.style.name if para.style else ""
        text = para.text.strip()
        if not text:
            continue

        if style.startswith("Heading"):
            flush()
            try:
                level = int(style.split()[-1])
            except ValueError:
                level = 1
            heading_stack = heading_stack[:level - 1] + [text]
            current_section = " / ".join(heading_stack)
        else:
            current_lines.append(text)

    # Extract tables
    for table in doc.tables:
        has_table = True
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        current_lines.append("[TABLE]\n" + "\n".join(rows))

    flush()
    return sections


def _extract_xlsx(file_path: str) -> list[PageSection]:
    import openpyxl
    wb = openpyxl.load_workbook(file_path, data_only=True)
    sections = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                rows.append(" | ".join(str(c) if c is not None else "" for c in row))

        if rows:
            sections.append(PageSection(
                text="\n".join(rows),
                section=sheet_name,
                sheet_name=sheet_name,
                has_table=True,
            ))

    return sections


def _extract_pptx(file_path: str) -> list[PageSection]:
    from pptx import Presentation
    prs = Presentation(file_path)
    sections = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
            if hasattr(shape, "notes_slide") and shape.notes_slide:
                notes_text = shape.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    lines.append(f"[Notes] {notes_text}")

        if lines:
            sections.append(PageSection(
                text="\n".join(lines),
                section=f"Slide {slide_num}",
                page_number=slide_num,
            ))

    return sections


def _extract_text(file_path: str) -> list[PageSection]:
    with open(file_path, encoding="utf-8", errors="replace") as f:
        text = f.read()

    # Split markdown by headings for better sections
    if file_path.endswith(".md"):
        return _split_markdown(text)

    return [PageSection(text=text.strip())] if text.strip() else []


def _split_markdown(text: str) -> list[PageSection]:
    import re
    sections = []
    heading_stack: list[str] = []
    current_lines: list[str] = []
    current_section = ""

    def flush():
        content = "\n".join(current_lines).strip()
        if content:
            sections.append(PageSection(text=content, section=current_section))
        current_lines.clear()

    for line in text.splitlines():
        heading_match = re.match(r"^(#{1,6})\s+(.*)", line)
        if heading_match:
            flush()
            level = len(heading_match.group(1))
            heading_text = heading_match.group(2).strip()
            heading_stack = heading_stack[:level - 1] + [heading_text]
            current_section = " / ".join(heading_stack)
        else:
            current_lines.append(line)

    flush()
    return sections


def _enrich_heading_paths(sections: list[PageSection]) -> None:
    """
    Infer heading paths from PDF text by detecting lines that look like section headers
    (all-caps short lines, or lines ending with no punctuation followed by a line break).
    Best-effort; not perfect for all PDFs.
    """
    import re
    heading_pattern = re.compile(r"^[A-Z][A-Z\s\d\-:]{3,60}$")
    current_heading = ""

    for sec in sections:
        lines = sec.text.splitlines()
        for line in lines[:3]:  # Only check first few lines of each page
            stripped = line.strip()
            if heading_pattern.match(stripped):
                current_heading = stripped.title()
                break
        if not sec.section and current_heading:
            sec.section = current_heading
